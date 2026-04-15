#!/usr/bin/env python3
"""
Generate a 64-slide PPTX presentation: Crypto Payment Services Deep Comparison
Services: Base Pay, Binance Pay, Coinbase Commerce, Stripe Crypto
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ============================================================
# 1. DESIGN SYSTEM CONSTANTS
# ============================================================

# Backgrounds
BG_PRIMARY = "#0F1923"
BG_SECONDARY = "#162033"
BG_CARD = "#1C2940"
BG_TABLE_HEADER = "#2D2D44"
BG_TABLE_ROW_ODD = "#182338"
BG_TABLE_ROW_EVEN = "#1E2D45"
BG_TITLE_BAR = "#0A1220"

# Text
TEXT_PRIMARY = "#FFFFFF"
TEXT_SECONDARY = "#B8C4D4"
TEXT_TERTIARY = "#7A8B9E"
TEXT_DARK = "#0F1923"

# Service accents
ACCENT_BASEPAY = "#0052FF"
ACCENT_BINANCE = "#F0B90B"
ACCENT_COINBASE = "#1652F0"
ACCENT_STRIPE = "#635BFF"

# Semantic
COLOR_POSITIVE = "#00C853"
COLOR_NEGATIVE = "#FF5252"
COLOR_NEUTRAL = "#FFD740"
COLOR_OPPORTUNITY = "#448AFF"
COLOR_THREAT = "#FF9100"

# Priority
COLOR_P0 = "#FF5252"
COLOR_P1 = "#FF9100"
COLOR_P2 = "#FFD740"

# Section colors
SECTION_PART0 = "#78909C"
SECTION_PART1A = "#0052FF"
SECTION_PART1B = "#F0B90B"
SECTION_PART1C = "#1652F0"
SECTION_PART1D = "#635BFF"
SECTION_PART2 = "#26C6DA"
SECTION_PART3 = "#AB47BC"

# Comparison table header bg per service
COMP_HDR_BASEPAY = "#1A2E5A"
COMP_HDR_BINANCE = "#3D3520"
COMP_HDR_COINBASE = "#1E2B55"
COMP_HDR_STRIPE = "#2A2655"

# Font
FONT_NAME = "Malgun Gothic"
FONT_FALLBACK = "Arial"

# Slide dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_str):
    """Convert hex color string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color_hex):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_shape(slide, left, top, width, height, fill_hex=None, border_hex=None, border_width=Pt(0)):
    """Add a rectangle shape with optional fill and border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.line.fill.background()
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        shape.fill.background()
    if border_hex:
        shape.line.fill.solid()
        shape.line.color.rgb = hex_to_rgb(border_hex)
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_hex=None, border_hex=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        shape.fill.background()
    if border_hex:
        shape.line.color.rgb = hex_to_rgb(border_hex)
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def set_text(tf, text, font_size=14, color=TEXT_SECONDARY, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Set text in a text frame."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(color)
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_paragraph(tf, text, font_size=14, color=TEXT_SECONDARY, bold=False, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(2), space_after=Pt(2), font_name=FONT_NAME, level=0):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(color)
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_textbox(slide, left, top, width, height, text="", font_size=14, color=TEXT_SECONDARY,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_text(tf, text, font_size, color, bold, alignment, font_name)
    return txBox


def add_section_indicator(slide, section_color):
    """Add thin top bar for section indicator."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04), fill_hex=section_color)


def add_title_bar(slide, title, section_color=None, font_size=26):
    """Add title bar to slide."""
    add_shape(slide, Inches(0), Inches(0.15), SLIDE_WIDTH, Inches(0.85), fill_hex=BG_TITLE_BAR)
    if section_color:
        add_shape(slide, Inches(0), Inches(0.15), Inches(0.08), Inches(0.85), fill_hex=section_color)
    add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12.133), Inches(0.75),
                title, font_size=font_size, color=TEXT_PRIMARY, bold=True)


def add_slide_number(slide, num):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
                str(num), font_size=9, color=TEXT_TERTIARY, alignment=PP_ALIGN.RIGHT)


def add_confidential(slide):
    """Add confidential footer."""
    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(3.0), Inches(0.3),
                "내부 공유용 / Confidential", font_size=8, color=TEXT_TERTIARY)


def add_source(slide, text, y=6.85):
    """Add source footnote."""
    add_textbox(slide, Inches(0.6), Inches(y), Inches(12.133), Inches(0.25),
                text, font_size=9, color=TEXT_TERTIARY)


def style_table_cell(cell, text, font_size=10, color=TEXT_SECONDARY, bold=False,
                      bg_hex=None, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Style a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(color)
    run.font.bold = bold
    run.font.name = font_name
    cell.text_frame.word_wrap = True
    # vertical centering
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # margins
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    if bg_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(bg_hex)


def add_styled_table(slide, rows, cols, data, col_widths, left=Inches(0.5), top=Inches(1.3),
                      header_bg=BG_TABLE_HEADER, row_height=Inches(0.38), font_size=10,
                      header_font_size=11, header_colors=None, width=None):
    """Add a styled table with headers and alternating rows."""
    if width is None:
        width = sum(col_widths)
    height = row_height * rows
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Style cells
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            if r == 0:
                # Header
                bg = header_bg
                if header_colors and c < len(header_colors) and header_colors[c]:
                    bg = header_colors[c]
                style_table_cell(cell, text, font_size=header_font_size, color=TEXT_PRIMARY,
                                  bold=True, bg_hex=bg, alignment=PP_ALIGN.CENTER)
            else:
                bg = BG_TABLE_ROW_ODD if r % 2 == 1 else BG_TABLE_ROW_EVEN
                style_table_cell(cell, text, font_size=font_size, color=TEXT_SECONDARY,
                                  bg_hex=bg)
            # Remove borders - set to subtle
            for border_name in ['top', 'bottom', 'left', 'right']:
                try:
                    border = getattr(cell, border_name, None)
                except:
                    pass

    return table_shape


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=14,
                        color=TEXT_SECONDARY, bold=False, bullet_char="\u2022 "):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = bullet_char + bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(color)
        run.font.bold = bold
        run.font.name = FONT_NAME
    return txBox


def create_base_slide(prs, section_color, title, slide_num):
    """Create a base slide with background, section indicator, title bar, number, footer."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, BG_PRIMARY)
    add_section_indicator(slide, section_color)
    add_title_bar(slide, title, section_color)
    add_slide_number(slide, slide_num)
    add_confidential(slide)
    return slide


def add_card(slide, x, y, w, h, accent_color, label, content_lines, label_size=12, content_size=11):
    """Add a card with accent top bar, label, and content."""
    add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill_hex=BG_CARD)
    add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.06), fill_hex=accent_color)
    add_textbox(slide, Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35),
                label, font_size=label_size, color=accent_color, bold=True)
    if content_lines:
        add_bullet_textbox(slide, Inches(x + 0.15), Inches(y + 0.5), Inches(w - 0.3), Inches(h - 0.6),
                           content_lines, font_size=content_size, bullet_char="")


def add_swot_slide(slide, strengths, weaknesses, opportunities, threats, font_size=10):
    """Add SWOT 2x2 grid to a slide that already has title bar."""
    # Strengths - top left
    add_card(slide, 0.5, 1.3, 6.0, 2.7, COLOR_POSITIVE, "Strengths (강점)", strengths, content_size=font_size)
    # Weaknesses - top right
    add_card(slide, 6.8, 1.3, 6.0, 2.7, COLOR_NEGATIVE, "Weaknesses (약점)", weaknesses, content_size=font_size)
    # Opportunities - bottom left
    add_card(slide, 0.5, 4.2, 6.0, 2.7, COLOR_OPPORTUNITY, "Opportunities (기회)", opportunities, content_size=font_size)
    # Threats - bottom right
    add_card(slide, 6.8, 4.2, 6.0, 2.7, COLOR_THREAT, "Threats (위협)", threats, content_size=font_size)


def add_two_column(slide, left_title, left_items, right_title, right_items,
                    left_color=COLOR_POSITIVE, right_color=COLOR_NEGATIVE):
    """Add two-column comparison layout."""
    # Left
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.3), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.06), fill_hex=left_color)
    add_textbox(slide, Inches(0.65), Inches(1.4), Inches(5.7), Inches(0.35),
                left_title, font_size=14, color=left_color, bold=True)
    add_bullet_textbox(slide, Inches(0.65), Inches(1.85), Inches(5.7), Inches(4.6),
                       left_items, font_size=11, bullet_char="\u2022 ")
    # Right
    add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.3), fill_hex=BG_CARD)
    add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.06), fill_hex=right_color)
    add_textbox(slide, Inches(6.95), Inches(1.4), Inches(5.7), Inches(0.35),
                right_title, font_size=14, color=right_color, bold=True)
    add_bullet_textbox(slide, Inches(6.95), Inches(1.85), Inches(5.7), Inches(4.6),
                       right_items, font_size=11, bullet_char="\u2022 ")


def add_kpi_card(slide, x, y, w, h, number, label, accent_color):
    """Add a KPI card with large number and small label."""
    add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill_hex=BG_CARD)
    add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.04), fill_hex=accent_color)
    add_textbox(slide, Inches(x), Inches(y + 0.25), Inches(w), Inches(0.65),
                number, font_size=28, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(y + 0.95), Inches(w), Inches(0.4),
                label, font_size=10, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)


def add_flow_steps(slide, steps, y_start, accent_color, x_start=0.5, box_w=1.8, box_h=0.6, gap=0.15):
    """Add horizontal flow step boxes."""
    n = len(steps)
    total_w = n * box_w + (n - 1) * gap
    # If total exceeds available width, shrink
    avail = 12.333
    if total_w > avail:
        box_w = (avail - (n - 1) * gap) / n

    for i, step_text in enumerate(steps):
        x = x_start + i * (box_w + gap)
        shape = add_rounded_rect(slide, Inches(x), Inches(y_start), Inches(box_w), Inches(box_h),
                                  fill_hex=BG_CARD, border_hex=accent_color, border_width=Pt(1))
        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, f"{i+1}. {step_text}", font_size=9, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)
        # Arrow between boxes (except last)
        if i < n - 1:
            arrow_x = x + box_w
            add_textbox(slide, Inches(arrow_x), Inches(y_start + box_h / 2 - 0.1),
                        Inches(gap), Inches(0.2), "\u2192", font_size=12, color=TEXT_TERTIARY,
                        alignment=PP_ALIGN.CENTER)


def add_numbered_blocks(slide, items, accent_color, y_start=1.3, block_h=0.85, gap=0.1):
    """Add vertical numbered insight blocks."""
    for i, (title, desc) in enumerate(items):
        y = y_start + i * (block_h + gap)
        add_shape(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(block_h), fill_hex=BG_CARD)
        add_shape(slide, Inches(0.6), Inches(y), Inches(0.04), Inches(block_h), fill_hex=accent_color)
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.17),
                                        Inches(0.4), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = hex_to_rgb(accent_color)
        badge.line.fill.background()
        set_text(badge.text_frame, str(i + 1), font_size=16, color=TEXT_PRIMARY,
                 bold=True, alignment=PP_ALIGN.CENTER)
        # Title + desc
        add_textbox(slide, Inches(1.3), Inches(y + 0.08), Inches(11.2), Inches(0.35),
                    title, font_size=12, color=TEXT_PRIMARY, bold=True)
        add_textbox(slide, Inches(1.3), Inches(y + 0.42), Inches(11.2), Inches(0.4),
                    desc, font_size=10, color=TEXT_SECONDARY)


def add_comparison_table(slide, headers, data, col_widths=None, top=Inches(1.3),
                          service_header_colors=None, font_size=9, header_font_size=10):
    """Add 4-service comparison table with colored headers."""
    rows = len(data) + 1
    cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(2.3)] + [Inches(2.5)] * (cols - 1)
    if service_header_colors is None:
        service_header_colors = [BG_TABLE_HEADER, COMP_HDR_BASEPAY, COMP_HDR_BINANCE,
                                   COMP_HDR_COINBASE, COMP_HDR_STRIPE]
    all_data = [headers] + data
    return add_styled_table(slide, rows, cols, all_data, col_widths, top=top,
                             header_colors=service_header_colors, font_size=font_size,
                             header_font_size=header_font_size)


# ============================================================
# 2. PRESENTATION CREATION
# ============================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ----------------------------------------------------------
    # SLIDE 1: Cover
    # ----------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_PRIMARY)
    add_section_indicator(slide, SECTION_PART0)

    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(1.2),
                "크립토 결제 서비스 4종 심층 비교 분석",
                font_size=36, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(0.8),
                "Base Pay / Binance Pay / Coinbase Commerce / Stripe Crypto",
                font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # 4-color accent line
    seg_w = 1.833
    colors = [ACCENT_BASEPAY, ACCENT_BINANCE, ACCENT_COINBASE, ACCENT_STRIPE]
    for i, c in enumerate(colors):
        add_shape(slide, Inches(3.0 + i * seg_w), Inches(4.3), Inches(seg_w), Inches(0.06), fill_hex=c)

    txb = add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10.333), Inches(0.5),
                      "결제(Payment) - 정산(Settlement) - 환불(Refund) 전 과정 비교",
                      font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.333), Inches(0.4),
                "작성일: 2026-04-15", font_size=12, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.7), Inches(10.333), Inches(0.4),
                "내부 공유용 / Confidential", font_size=12, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, 1)

    # ----------------------------------------------------------
    # SLIDE 2: TOC
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "목차", 2)

    toc_items = [
        (SECTION_PART0, "Part 0: 도입부", "#1~#8 -- Executive Summary, 시장 현황"),
        (SECTION_PART1A, "Part 1-A: Base Pay", "#9~#20 -- 서비스 개요, 결제/정산/환불, SWOT"),
        (SECTION_PART1B, "Part 1-B: Binance Pay", "#21~#28 -- 서비스 개요, 결제/정산/환불, SWOT"),
        (SECTION_PART1C, "Part 1-C: Coinbase Commerce", "#29~#36 -- 서비스 개요, 결제/정산/환불, SWOT"),
        (SECTION_PART1D, "Part 1-D: Stripe Crypto", "#37~#48 -- 서비스 개요, 결제/정산/환불, SWOT"),
        (SECTION_PART2, "Part 2: 교차 비교 분석", "#49~#56 -- 수수료, 정산, 환불, 기능, 스코어카드"),
        (SECTION_PART3, "Part 3: 전략적 시사점 및 결론", "#57~#64 -- 트렌드, 권고, 결론"),
    ]
    # Divider line
    add_shape(slide, Inches(4.8), Inches(1.3), Inches(0.04), Inches(5.5), fill_hex=TEXT_TERTIARY)

    for i, (color, part_label, desc) in enumerate(toc_items):
        y = 1.3 + i * 0.75
        # Colored marker
        add_shape(slide, Inches(0.6), Inches(y + 0.1), Inches(0.15), Inches(0.15), fill_hex=color)
        add_textbox(slide, Inches(0.9), Inches(y), Inches(3.7), Inches(0.65),
                    part_label, font_size=14, color=color, bold=True)
        add_textbox(slide, Inches(5.0), Inches(y), Inches(7.9), Inches(0.65),
                    desc, font_size=12, color=TEXT_SECONDARY)

    # ----------------------------------------------------------
    # SLIDE 3: Executive Summary - 4 Quad
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "Executive Summary -- 전체 핵심 인사이트", 3)

    quad_data = [
        (0.5, 1.3, ACCENT_BASEPAY, "Base Pay",
         "업계 유일 온체인 에스크로 기반 Authorize-Capture 프로토콜. Coinbase 생태계(6,000만+) + Shopify(550만) + x402 AI 결제(연간화 $600M). 약점: 법정화폐 직접 정산 미지원"),
        (6.8, 1.3, ACCENT_BINANCE, "Binance Pay",
         "누적 거래량 $280B+, 2,000만 가맹점 세계 최대. 오프체인 10,000+ TPS 즉시 정산. B2C 98% 스테이블코인. 약점: 법정화폐 정산 불가, 70개국+ 미진출"),
        (0.5, 4.3, ACCENT_COINBASE, "Coinbase Commerce",
         "시장점유율 12%(3위). Base L2 + Shopify 파트너십. 비용 우위(1% vs 2.9%). 약점: 환불 업계 최하위, 2026.03 Commerce 종료"),
        (6.8, 4.3, ACCENT_STRIPE, "Stripe Crypto",
         "Bridge($1.1B)+Tempo+Privy 수직 통합 '스테이블코인 풀스택 클라우드'. 가맹점 크립토 노출 완전 제거. 약점: 미국 전용, 1.5% 프리미엄"),
    ]
    for x, y, accent, label, content in quad_data:
        add_card(slide, x, y, 6.0, 2.8, accent, label, [content])

    # ----------------------------------------------------------
    # SLIDE 4: Executive Summary - Table
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "Executive Summary -- 서비스별 한 줄 요약", 4)
    data = [
        ["서비스", "핵심 포지션", "최대 강점", "최대 약점"],
        ["Base Pay", "프로그래머블 온체인 결제 인프라", "에스크로, 오픈소스, 다층 수익", "법정화폐 정산 미지원, 고객 지원 부재"],
        ["Binance Pay", "거래소 생태계 내장형 대규모 결제", "3억 사용자, 즉시정산, 가스비 0", "법정화폐 정산 불가, 폐쇄적 생태계"],
        ["Coinbase Commerce", "기관 통합형 온체인 결제 게이트웨이", "Base L2 소유, Shopify, 비용 우위", "환불 자동화 부재, 플랫폼 전환 불확실성"],
        ["Stripe Crypto", "전통 결제 + 크립토 수직 통합", "zero-crypto 경험, 개발자 UX", "미국 전용, 1.5% 프리미엄, T+2 정산"],
    ]
    col_w = [Inches(1.8), Inches(3.5), Inches(3.5), Inches(3.5)]
    add_styled_table(slide, 5, 4, data, col_w, font_size=11, header_font_size=12)

    # ----------------------------------------------------------
    # SLIDE 5: Market Overview - KPI Cards
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "크립토 결제 시장 현황 개요", 5)
    add_kpi_card(slide, 0.5, 1.3, 2.9, 1.5, "$2.95B", "결제 앱 시장 (2035)", SECTION_PART2)
    add_kpi_card(slide, 3.65, 1.3, 2.9, 1.5, "$4.74B", "게이트웨이 시장 (2030)", COLOR_OPPORTUNITY)
    add_kpi_card(slide, 6.8, 1.3, 2.9, 1.5, "$89.4B", "스테이블코인 인프라 (2034)", SECTION_PART3)
    add_kpi_card(slide, 9.95, 1.3, 2.9, 1.5, "+43%", "미국 결제 사용 증가", COLOR_POSITIVE)

    bullets = [
        "스테이블코인 시가총액: ~$3,086억 (2026.01) -> 연말 $1조 초과 전망",
        "스테이블코인 연간 결제 볼륨: $3,900억 (실질 결제, McKinsey)",
        "B2B 스테이블코인 거래량: $4,000억 (전년 대비 2배, PYMNTS)",
        "미국 가맹점 암호화폐 결제 수용률: 39%",
    ]
    add_bullet_textbox(slide, Inches(0.6), Inches(3.1), Inches(12.1), Inches(3.5),
                       bullets, font_size=12)
    add_source(slide, "출처: Research Nester, GII Research, McKinsey, PYMNTS, CoinLaw")

    # ----------------------------------------------------------
    # SLIDE 6: 5 Key Trends
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "크립토 결제 시장 -- 5대 핵심 트렌드", 6)
    trends = [
        ("스테이블코인 결제 폭발적 성장", "2025년 $33T 거래량, Binance Pay B2C의 98%가 스테이블코인"),
        ("L2 기반 결제 실용화", "가스비 $0.01 이하 + 서브초 정산으로 상거래 수준 도달"),
        ("AI 에이전트 결제(Machine Payments) 부상", "x402, MPP, ACP 등 프로토콜 경쟁 시작"),
        ("전통 결제사의 크립토 진출", "Stripe(1.5%), PayPal(PYUSD), Visa가 본격 참여"),
        ("규제 명확화", "US GENIUS Act(2027.01 발효), EU MiCA 전면 시행"),
    ]
    add_numbered_blocks(slide, trends, SECTION_PART2, y_start=1.3, block_h=0.95, gap=0.12)

    # ----------------------------------------------------------
    # SLIDE 7: 3 Camps
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "크립토 결제 시장 경쟁 구도 -- 3대 진영", 7)
    camps = [
        (0.5, ACCENT_BINANCE, "진영 1: 거래소 내장형",
         ["Binance Pay, Crypto.com Pay", "자사 사용자 기반", "크립토-크립토 정산"]),
        (4.7, SECTION_PART0, "진영 2: 전통 게이트웨이형",
         ["BitPay, CoinGate, NOWPayments", "크립토 수취 -> 법정화폐 전환"]),
        (8.9, ACCENT_STRIPE, "진영 3: 전통 결제 확장형",
         ["PayPal, Stripe, Coinbase Commerce", "기존 인프라에 크립토 통합"]),
    ]
    for x, accent, label, items in camps:
        add_card(slide, x, 1.3, 3.9, 4.5, accent, label, items, content_size=11)

    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                "시장 점유율: BitPay 20% | CoinGate 14% | Coinbase Commerce 12% | Binance Pay 8%",
                font_size=11, color=TEXT_TERTIARY)

    # ----------------------------------------------------------
    # SLIDE 8: Positioning Map
    # ----------------------------------------------------------
    slide = create_base_slide(prs, SECTION_PART0, "경쟁사 포지셔닝 맵", 8)

    # Draw axes
    add_shape(slide, Inches(1.5), Inches(3.8), Inches(10.5), Inches(0.02), fill_hex=TEXT_TERTIARY)  # X
    add_shape(slide, Inches(6.75), Inches(1.3), Inches(0.02), Inches(5.2), fill_hex=TEXT_TERTIARY)  # Y
    add_textbox(slide, Inches(5.0), Inches(6.55), Inches(5.0), Inches(0.3),
                "법정화폐 통합도 →", font_size=10, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.3), Inches(3.3), Inches(1.2), Inches(0.3),
                "온체인 ↑", font_size=10, color=TEXT_TERTIARY)

    # Quadrant labels
    add_textbox(slide, Inches(1.5), Inches(1.3), Inches(4.5), Inches(0.3),
                "프로그래머블 + 크립토 네이티브", font_size=9, color=TEXT_TERTIARY)
    add_textbox(slide, Inches(7.5), Inches(1.3), Inches(4.5), Inches(0.3),
                "프로그래머블 + 법정화폐 통합 (최대 위협)", font_size=9, color=COLOR_NEGATIVE)
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(4.5), Inches(0.3),
                "단순 + 크립토 네이티브", font_size=9, color=TEXT_TERTIARY)
    add_textbox(slide, Inches(7.5), Inches(4.0), Inches(4.5), Inches(0.3),
                "단순 + 법정화폐 통합", font_size=9, color=TEXT_TERTIARY)

    # Bubbles
    bubbles = [
        (3.0, 2.0, 0.7, ACCENT_BASEPAY, "Base Pay"),
        (9.5, 2.0, 0.8, ACCENT_STRIPE, "Stripe"),
        (3.5, 5.0, 0.9, ACCENT_BINANCE, "Binance Pay"),
        (6.0, 3.0, 0.6, ACCENT_COINBASE, "Coinbase"),
        (9.0, 5.0, 0.6, SECTION_PART0, "PayPal"),
    ]
    for bx, by, br, bc, bname in bubbles:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx), Inches(by), Inches(br), Inches(br))
        oval.fill.solid()
        oval.fill.fore_color.rgb = hex_to_rgb(bc)
        oval.line.fill.background()
        set_text(oval.text_frame, bname, font_size=9, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    # ==================================================================
    # PART 1-A: BASE PAY (Slides 9-20)
    # ==================================================================

    # SLIDE 9
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 서비스 개요", 9)
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.0),
                "Coinbase L2 블록체인 Base 네트워크 위에 구축된 결제 인프라 생태계.\n\n"
                "핵심: Commerce Payments Protocol -- 오픈소스 온체인 Authorize-Capture 결제 프로토콜.\n\n"
                "업계 유일의 온체인 에스크로 기반 결제 시스템으로, 전통 결제의 승인-캡처-환불-취소 흐름을 스마트 컨트랙트로 재현.",
                font_size=12, color=TEXT_SECONDARY)
    data = [
        ["항목", "내용"],
        ["정산 속도", "Base ~2초"],
        ["거래 수수료", "가스비 ~$0.01 + 플랫폼 1%"],
        ["결제 보장", "머천트 정확한 요청 금액 수령"],
        ["토큰 유연성", "Uniswap V3 유동성 있는 모든 토큰"],
        ["오픈소스", "Apache 2.0 라이선스"],
    ]
    add_styled_table(slide, 6, 2, data, [Inches(2.0), Inches(4.5)], left=Inches(6.3), top=Inches(1.3))

    # SLIDE 10
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 제품/프로토콜 라인업 및 타임라인", 10)
    data = [
        ["프로토콜", "세대", "핵심 기능", "배포 상태"],
        ["Transfers.sol (1세대)", "즉시 결제", "직접 토큰 전송, Uniswap V3 스왑", "Base, ETH, Polygon"],
        ["AuthCaptureEscrow (2세대)", "에스크로 결제", "Authorize-Capture-Void-Refund", "Base"],
        ["x402 프로토콜", "머신 결제", "HTTP 402 기반 AI 에이전트 결제", "Base, Solana"],
        ["Coinbase Business", "통합 플랫폼", "커스터디, 법정화폐 정산", "미국/싱가포르"],
    ]
    add_styled_table(slide, 5, 4, data, [Inches(2.8), Inches(1.5), Inches(4.5), Inches(3.5)],
                      top=Inches(1.3), font_size=10)
    # Timeline
    milestones = ["~2024\nCommerce", "2025중반\nProtocol 공개", "2025.09\nx402 출시", "2026.02\nStripe x402", "2026.03.31\nCommerce 종료"]
    add_shape(slide, Inches(1.0), Inches(5.5), Inches(11.0), Inches(0.02), fill_hex=TEXT_TERTIARY)
    for i, ms in enumerate(milestones):
        x = 1.0 + i * 2.75
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.4), Inches(0.2), Inches(0.2))
        oval.fill.solid()
        c = COLOR_NEGATIVE if "종료" in ms else ACCENT_BASEPAY
        oval.fill.fore_color.rgb = hex_to_rgb(c)
        oval.line.fill.background()
        add_textbox(slide, Inches(x - 0.5), Inches(5.65), Inches(1.5), Inches(0.7),
                    ms, font_size=8, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

    # SLIDE 11
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 결제 시나리오 A: 즉시 결제 (Charge)", 11)
    steps = ["결제 요청 생성", "TransferIntent\n생성", "지갑 연결\n결제 승인", "스마트 컨트랙트\n실행(자동 스왑)", "머천트 지갑\n정산 전송", "이벤트 발행\n(~2초)"]
    add_flow_steps(slide, steps, 1.4, ACCENT_BASEPAY, box_w=1.85, box_h=0.7)
    data = [
        ["항목", "금액", "부담 주체"],
        ["플랫폼 수수료", "1%", "머천트 (정산에서 차감)"],
        ["네트워크 가스비 (Base)", "~$0.01", "구매자"],
        ["DEX 스왑 슬리피지", "변동", "구매자 (토큰 다를 경우)"],
    ]
    add_styled_table(slide, 4, 3, data, [Inches(3.0), Inches(4.0), Inches(5.3)], top=Inches(4.2))
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.4),
                "적합 사용처: 디지털 상품, 구독 활성화, 에스크로 불필요 시나리오",
                font_size=11, color=TEXT_TERTIARY)

    # SLIDE 12
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 결제 시나리오 B: 인가-캡처 결제", 12)
    steps = ["PaymentInfo\n생성", "구매자\n결제 승인", "Authorize\n에스크로 예치", "[대기 기간]\n주문/배송 확인", "Capture\n머천트 전송", "미캡처분\n자동 반환"]
    add_flow_steps(slide, steps, 1.4, ACCENT_BASEPAY, box_w=1.85, box_h=0.7)
    bullets = [
        "에스크로 메커니즘: Minimal Proxy 패턴, 오퍼레이터별 격리, 비업그레이드, ReentrancyGuard",
        "부분 캡처 지원 -- 분할 배송에도 유연하게 대응",
        "수수료: 1% (Capture 시 차감), 가스비 ~$0.01 (Authorize/Capture 각각), Void 시 수수료 없음",
        "적합 사용처: 실물 이커머스, 예약/호텔, 세금 확정 후 결제",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.0), bullets, font_size=11)

    # SLIDE 13
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 결제 시나리오 C: x402 머신/AI 에이전트 결제", 13)
    steps = ["AI 에이전트\nHTTP GET", "서버 HTTP 402\n+ Invoice", "클라이언트\n결제 서명", "PAYMENT-SIG\n헤더 재요청", "Facilitator\n검증+시뮬레이션", "온체인 정산\nHTTP 200"]
    add_flow_steps(slide, steps, 1.4, ACCENT_BASEPAY, box_w=1.85, box_h=0.7)
    data = [
        ["스킴", "설명", "환불"],
        ["exact", "고정 금액 결제", "별도 프로세스 필요"],
        ["upto", "사용량 기반 최대 금액 설정", "미사용분 미과금"],
    ]
    add_styled_table(slide, 3, 3, data, [Inches(2.0), Inches(5.5), Inches(4.8)], top=Inches(4.0))
    add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.8),
                "성장: Base 1.19억+ 트랜잭션 | 연간화 ~$600M | USDC 98.7% | Foundation: Google, AWS, Visa, Circle, Anthropic",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 14
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 지원 네트워크/토큰 및 통합 방식", 14)
    data = [
        ["체인", "환경", "1세대 컨트랙트", "2세대 컨트랙트"],
        ["Base", "Mainnet", "지원", "지원"],
        ["Ethereum", "Mainnet", "지원", "--"],
        ["Polygon", "Mainnet", "지원", "--"],
    ]
    add_styled_table(slide, 4, 4, data, [Inches(2.5), Inches(2.5), Inches(3.5), Inches(3.8)], top=Inches(1.3))
    data2 = [
        ["지갑 유형", "토큰", "권장 Collector", "UX 품질"],
        ["EOA", "USDC/EURC", "ERC3009PaymentCollector", "최상 (1 서명)"],
        ["EOA", "기타 ERC-20", "Permit2PaymentCollector", "양호"],
        ["Coinbase Smart Wallet", "모든 토큰", "SpendPermissionCollector", "최상 (구독)"],
    ]
    add_styled_table(slide, 4, 4, data2, [Inches(2.8), Inches(2.5), Inches(3.8), Inches(3.2)], top=Inches(3.5))
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
                "통합: Shopify 플러그인(매우 낮음) | Commerce API(낮음) | 스마트 컨트랙트(높음) | x402 HTTP(낮음)",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 15
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 정산 시나리오", 15)
    bullets = [
        "암호화폐 직접 정산 -- 즉시 정산(Charge): ~2초, 에스크로 정산(Capture): Capture 시 즉시",
        "정산 통화: USDC(기본), ETH/기타, 자동 변환(Uniswap V3 DEX)",
        "법정화폐 정산 경로 (간접):",
        "  1) Coinbase Business 경유 (미국/싱가포르만): USDC -> fiat -> 은행 (+1-5 영업일)",
        "  2) Coinbase 거래소 경유 (글로벌): USDC -> USD -> 은행 출금",
        "  3) Shopify 통합 경유: Stripe 인프라 활용 현지 통화 정산 (수수료 0%)",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.8), bullets, font_size=12)
    data = [
        ["시나리오", "머천트 수령액", "구매자 총비용"],
        ["USDC -> USDC (스왑 없음)", "$99.00", "~$100.01"],
        ["ETH -> USDC (스왑 있음)", "$99.00", "~$100.11~$100.51"],
    ]
    add_styled_table(slide, 3, 3, data, [Inches(4.0), Inches(4.0), Inches(4.3)], top=Inches(4.5))

    # SLIDE 16
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 환불 시나리오 (4가지)", 16)
    data = [
        ["시나리오", "설명", "수수료", "속도"],
        ["A. Void (캡처 전 취소)", "전체 잔여 금액 즉시 반환", "미발생 (~$0.01)", "~2초"],
        ["B. Refund (캡처 후)", "RefundCollector로 부분/전액", "가스비 ~$0.01", "~2초"],
        ["C. Reclaim (만료 후)", "구매자 직접 reclaim() 호출", "가스비 ~$0.01", "~2초"],
        ["D. Charge 후 환불", "1세대: 수동 / 2세대: refund()", "가스비 ~$0.01", "~2초"],
    ]
    add_styled_table(slide, 5, 4, data, [Inches(2.5), Inches(5.0), Inches(2.5), Inches(2.3)], top=Inches(1.3))
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.2),
                "refundExpiry: 결제 생성 시 사전 정의, 만료 후 환불 불가\n"
                "분쟁 해결: 차지백 없음 (머천트 유리), 에스크로가 부분적 보완\n"
                "상태 머신: Created -> Authorized -> Captured/Voided/Reclaimed -> Refunded",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 17
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 수익 구조 -- 5개 레이어", 17)
    layers = [
        (11.0, "#0A2E6B", "Layer 0: USDC Revenue Share -- $332.5M/분기 (Coinbase 매출의 22%)"),
        (9.5, "#0C3577", "Layer 1: Base 시퀀서 수수료 -- $75.4M/년 (L2 전체의 62%)"),
        (8.0, "#0E3C83", "Layer 2: x402 Facilitator 수수료 -- 월 1,000건 무료, 이후 과금"),
        (6.5, "#104390", "Layer 3: Commerce 결제 수수료 -- 1% 고정"),
        (5.0, "#12499D", "Layer 4: Coinbase Business -- 커스터디, 환전, 출금 수수료"),
    ]
    for i, (w, color, text) in enumerate(layers):
        x = (13.333 - w) / 2
        y = 1.4 + i * 1.0
        shape = add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.85), fill_hex=color)
        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, text, font_size=11, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                "전략: Commerce 1%는 박한 마진이나, Base 트랜잭션 볼륨 + USDC 유통량 + 사용자 온보딩을 동시 달성하는 인프라",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 18
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 기술 아키텍처", 18)
    # Two columns
    add_card(slide, 0.5, 1.3, 5.9, 3.5, ACCENT_BASEPAY, "1세대 Transfers.sol",
             ["즉시 결제, 9가지 함수", "Permit2 + Uniswap V3", "3개 체인 배포 (Base, ETH, Polygon)"])
    add_card(slide, 6.7, 1.3, 5.9, 3.5, ACCENT_BASEPAY, "2세대 AuthCaptureEscrow",
             ["에스크로 기반 6가지 연산", "Token Collector 모듈 5종", "PaymentInfo 구조체, Base 전용"])
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5),
                "기술 스택: Solidity ^0.8.17, Foundry, OpenZeppelin, Uniswap V3, Apache 2.0\n"
                "보안: Spearbit + Coinbase Protocol Security 감사, 비업그레이드 컨트랙트, ReentrancyGuard\n"
                "x402: TypeScript/Python/Rust SDK, EIP-3009, Witness 패턴, 프라이빗 멤풀",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 19
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] 사용자 인사이트 -- 강점 & 약점", 19)
    add_two_column(slide,
                   "호평 / Strengths",
                   ["서브초 정산 속도", "낮은 수수료 (카드 대비 66% 절감)", "Shopify 네이티브 통합", "차지백 리스크 제거", "에스크로 (Auth-Capture)"],
                   "Pain Points / 약점",
                   ["[심각] 고객 지원 부재 (4-5일 응답)", "[심각] Commerce 종료/강제 마이그레이션", "[높음] 법정화폐 직접 정산 미지원", "[높음] 수동 환불 프로세스", "[높음] 자금 인출 어려움"])
    # Bottom callout
    add_shape(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(6.7), Inches(0.04), Inches(0.4), fill_hex=COLOR_NEGATIVE)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                "시드 구문 보안 인시던트 (2026.03) | 평점: Capterra 4.4/5, BBB F",
                font_size=10, color=COLOR_NEGATIVE)

    # SLIDE 20
    slide = create_base_slide(prs, SECTION_PART1A, "[Base Pay] SWOT 분석", 20)
    add_swot_slide(slide,
                   ["업계 유일 온체인 에스크로 (Auth-Capture)", "다층 수익 구조 (시퀀서비+USDC+Commerce+x402)", "Coinbase 생태계 (6,000만+)", "1% 수수료, 서브초 정산, 오픈소스", "규제 친화적 (미국 상장사, MSB)"],
                   ["법정화폐 직접 정산 미지원", "구매자 보호 부재 (차지백 없음)", "고객 지원 부재 (BBB F)", "글로벌 접근성 제한 (미국/싱가포르만)", "환불 UX 미흡, 비업그레이드 리스크"],
                   ["AI 에이전트 결제 (x402, 연간화 $600M)", "규제 명확화 (GENIUS Act 2027.01)", "Shopify 550만 가맹점 미개척", "법정화폐+온체인 공백 영역"],
                   ["Stripe Stablecoin Shopify 롤아웃 [높음]", "Solana Pay Visa/Mastercard 파트너십", "PayPal PYUSD 70개 시장 확장", "Commerce 종료 신뢰 훼손 [높음]"])

    # ==================================================================
    # PART 1-B: BINANCE PAY (Slides 21-28)
    # ==================================================================

    # SLIDE 21
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 서비스 개요 및 핵심 수치", 21)
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.0),
                "Binance가 2021년 출시한 오프체인 암호화폐 결제 서비스.\n\n"
                "P2P 송금(무료) + 가맹점 결제(API/Checkout/Payment Links) + Payout(배치) 통합.\n\n"
                "세계 최대 암호화폐 거래소의 결제 인프라로, 3억+ 사용자 기반 활용.",
                font_size=12, color=TEXT_SECONDARY)
    data = [
        ["지표", "수치"],
        ["누적 거래량", "$280B+"],
        ["연간 거래량", "$121B (2025)"],
        ["누적 거래 건수", "13.6억 건"],
        ["Pay 수수료 매출", "~$110M (2025)"],
        ["Binance 등록 사용자", "3억 명+"],
        ["가맹점 수", "2,000만+ *"],
        ["P2P 활성 사용자", "4,500만 명+"],
    ]
    add_styled_table(slide, 8, 2, data, [Inches(2.5), Inches(4.0)], left=Inches(6.3), top=Inches(1.3))

    # SLIDE 22
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 결제 시나리오 (4가지)", 22)
    panels = [
        (0.5, 1.3, "시나리오 A: P2P 개인 간 결제", ["Pay ID/QR/이메일/전화번호", "수수료 0%, 즉시, 300종+"]),
        (6.7, 1.3, "시나리오 B: Merchant Hosted Checkout", ["Create Order V2 API", "50종+ 결제, 자동 환전, 즉시"]),
        (0.5, 4.1, "시나리오 C: Native API / Payment Links", ["가맹점 자체 UI QR 또는 딥링크", "코딩 불필요 링크 생성"]),
        (6.7, 4.1, "시나리오 D: QR 코드 오프라인 결제", ["API 주문 생성 -> QR 표시", "Binance 앱 스캔 -> 즉시 처리"]),
    ]
    for x, y, label, items in panels:
        add_card(slide, x, y, 5.9, 2.5, ACCENT_BINANCE, label, items, content_size=10)
    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
                "전제조건: 송수신 양측 Binance 계정 + KYC 인증 필수",
                font_size=10, color=COLOR_NEGATIVE)

    # SLIDE 23
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 수수료 구조", 23)
    data = [
        ["수수료 항목", "요율", "비고"],
        ["P2P 송금", "0%", "무료"],
        ["가맹점 결제 (MDR)", "1.0%", "고객 결제 금액 기준"],
        ["Payout", "0.80% (최대 $5/건)", "2024.12부터 유료 전환"],
        ["블록체인 가스비", "0%", "오프체인"],
        ["FX 스프레드 (숨겨진 비용)", "추정 0.1~0.5%", "자동 환전 시 내재 마진, 공식 비공개"],
        ["P2P 분쟁 수수료", "4회차부터 부과", "처음 3회 무료"],
    ]
    add_styled_table(slide, 7, 3, data, [Inches(3.5), Inches(3.5), Inches(5.3)], top=Inches(1.3))
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.7), fill_hex=BG_CARD)
    add_textbox(slide, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.6),
                "$10,000 매출 법정화폐 최종 수취 시뮬레이션: 총비용 ~$140~150 (실효 ~1.4~1.5%), 수동 처리 필요",
                font_size=13, color=TEXT_PRIMARY, bold=True)

    # SLIDE 24
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 정산 시나리오 -- 오프체인 즉시 정산", 24)
    steps = ["고객 Funding\nWallet 차감", "FX Engine\nBTC->USDT 변환", "MDR 1%\n차감", "가맹점 Wallet\nUSDT 즉시 입금"]
    add_flow_steps(slide, steps, 1.4, ACCENT_BINANCE, box_w=2.8, box_h=0.7)
    data = [
        ["항목", "수치"],
        ["TPS", "10,000+ (4노드 클러스터)"],
        ["트랜잭션 지연", "대부분 10ms 이내"],
        ["페일오버 시간", "1초 이내"],
        ["아키텍처", "Raft 합의 + RocksDB + 관계형 DB"],
    ]
    add_styled_table(slide, 5, 2, data, [Inches(3.0), Inches(9.3)], top=Inches(3.5))
    add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(6.3), Inches(0.04), Inches(0.5), fill_hex=COLOR_NEGATIVE)
    add_textbox(slide, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.5),
                "법정화폐 전환 한계: 직접 은행 입금 불가 -- 거래소 수동 환전 필요 (추가 0.1% + 출금 수수료)",
                font_size=11, color=COLOR_NEGATIVE)

    # SLIDE 25
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 환불 시나리오", 25)
    bullets = [
        "전액 환불: Refund Order API, prepayId 기반, MDR도 비례 환불",
        "부분 환불: 여러 차례 가능, 누적 <= 원래 금액",
        "환불 통화: 정산 통화(USDT)로 환불 (역환전 없음) -- BTC 결제 시 USDT 환불 (가격 상승분 손실)",
        "환불 속도: 즉시 (오프체인)",
        "분쟁 해결: 차지백 없음, 가맹점 재량, Binance 중개 가능 (강제력 제한)",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0), bullets, font_size=12)
    data = [
        ["에러 코드", "설명"],
        ["400202", "ORDER_NOT_FOUND"],
        ["400607", "PAYMENT_REFUND_TOO_MANY_TIMES"],
        ["400608", "PAYMENT_REFUND_AMOUNT_INVALID"],
        ["400611", "PAYMENT_INSUFFICIENT_BALANCE"],
    ]
    add_styled_table(slide, 5, 2, data, [Inches(2.0), Inches(10.3)], top=Inches(4.5))

    # SLIDE 26
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 비즈니스 모델 및 기술 아키텍처", 26)
    data = [
        ["분류", "설명", "추정 매출"],
        ["1차: MDR 1%", "Merchant 결제 수수료", "~$110M"],
        ["2차: Payout 0.80%", "배치 출금 수수료", "미공개"],
        ["3차: FX 스프레드", "자동 환전 시 내재 마진", "미공개"],
        ["4차: 생태계 유입", "거래소 거래 수수료 간접 기여", "간접적"],
    ]
    add_styled_table(slide, 5, 3, data, [Inches(2.5), Inches(5.0), Inches(4.8)], top=Inches(1.3))
    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.5),
                "Binance 전체 내 비중: $110M / $17.5B = 0.6% -- '독립 수익 사업이 아닌 생태계 전략 도구'\n\n"
                "기술: Java/Spring Boot, RocksDB+Raft, Kafka/RabbitMQ, Redis, Docker/K8s\n"
                "보안: HMAC-SHA512 + 1초 윈도우 + 32자 Nonce, KYC 3단계, AML 이중 레이어\n"
                "컴플라이언스 누적 투자: $1.2B",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 27
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] 사용자 인사이트 -- 강점 & 약점", 27)
    add_two_column(slide,
                   "호평",
                   ["즉시 처리, P2P 무료", "300종+ 지원, QR 편의성", "앱 내 통합", "통합 속도, Hosted Checkout", "자동 환전 정산"],
                   "불만",
                   ["고객 지원 품질 (10/10)", "계정 동결/잠금 (9/10)", "Binance 계정 필수 (8/10)", "규제 지역 차단 (8/10)", "법정화폐 직접 정산 불가 [매우 높음]"])
    add_shape(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(6.7), Inches(0.04), Inches(0.4), fill_hex=COLOR_NEUTRAL)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
                "평점: Trustpilot 1.4/5 vs App Store 4.7/5 -- 극단적 괴리",
                font_size=10, color=COLOR_NEUTRAL)

    # SLIDE 28
    slide = create_base_slide(prs, SECTION_PART1B, "[Binance Pay] SWOT 분석", 28)
    add_swot_slide(slide,
                   ["3억 사용자, 2,000만 가맹점", "오프체인 즉시 정산 + 가스비 0 (10,000+ TPS)", "300종+ 암호화폐, P2P 무료, 1% MDR", "B2C 98% 스테이블코인, 생태계 시너지"],
                   ["법정화폐 직접 정산 불가 (최대 구조적 약점)", "70개국+ 미진출 (미국/영국 포함)", "Binance 계정 필수 (폐쇄적)", "구독 결제 미지원, 고객 지원 1.4/5"],
                   ["신흥국 금융 포용 수요 급증", "법정화폐 정산 추가 시 독보적 포지션", "B2B 국경 간 결제 시장", "CBDC 브릿지 역할"],
                   ["PayPal/Stripe 크립토 통합 (4억+)", "스테이블코인 규제 강화 (MiCA)", "CBDC 출시로 수요 잠식", "Coinbase 오픈소스 전략"])

    # ==================================================================
    # PART 1-C: COINBASE COMMERCE (Slides 29-36)
    # ==================================================================

    # SLIDE 29
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 서비스 개요", 29)
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5),
                "가맹점 암호화폐 결제 게이트웨이. 2018년 출시.\n"
                "스마트 컨트랙트 기반 온체인 프로토콜로 전환 중.\n"
                "시장점유율 12% (3위).",
                font_size=12, color=TEXT_SECONDARY)
    data = [
        ["시기", "제품", "특징"],
        ["2018-2024", "Commerce Legacy", "전통 API 기반"],
        ["2024-2025", "Commerce Onchain", "스마트 컨트랙트 기반"],
        ["2025.06", "Coinbase Payments", "Shopify 연동, Base 기반"],
        ["2025.10", "Commerce Payments Protocol", "오픈소스 인가-캡처"],
        ["2026.01", "Commerce 업데이트", "Coinbase 유저 간 무료"],
        ["2026 진행", "Coinbase Business", "법정화폐 인출 강화"],
    ]
    add_styled_table(slide, 7, 3, data, [Inches(1.8), Inches(3.0), Inches(3.5)], left=Inches(6.3), top=Inches(1.3), font_size=9)
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(0.04), Inches(0.5), fill_hex=COLOR_NEGATIVE)
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.5),
                "핵심 변곡점: 2026.03.31 Commerce 종료. 미국/싱가포르만 Business 전환 가능. 기타 국가 강제 이탈",
                font_size=11, color=COLOR_NEGATIVE)

    # SLIDE 30
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 결제 시나리오", 30)
    add_card(slide, 0.5, 1.3, 6.0, 4.5, ACCENT_COINBASE, "기본 결제 (Charge)",
             ["1. Commerce API (POST /charges) -> Charge 객체", "2. hosted_url 접속 -> 토큰/네트워크 선택",
              "3. 온체인 트랜잭션 + DEX USDC 자동 환전", "4. 1% 수수료 차감 -> 가맹점 지갑 정산"])
    add_card(slide, 6.8, 1.3, 6.0, 4.5, COLOR_POSITIVE, "Protocol (Auth-Capture)",
             ["1. Authorization: 에스크로에 자금 예치", "2. Capture: 상품 준비 후 에스크로 -> 가맹점",
              "3. Void: 캡처 전 전액 반환", "4. Reclaim: 인가 만료 후 소비자 직접 회수"])
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
                "지원: BTC, ETH, LTC, BCH, DOGE, USDC, DAI, EURC 등 100+종 | Base (~200ms, ~$0.01), ETH, Polygon | API 10,000 req/hr",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 31
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 수수료 구조", 31)
    data = [
        ["결제 유형", "수수료", "비고"],
        ["Coinbase 유저 -> 가맹점", "무료, 즉시", "2026.01 업데이트"],
        ["외부 지갑 -> 가맹점", "1% + 가스비(~$0.01)", "Base 기준"],
        ["Shopify USDC", "1%", "외환 수수료 없음"],
        ["DEX 스왑 포함 시", "~1.3% + 가스비", "Uniswap V3 0.3% 추가"],
    ]
    add_styled_table(slide, 5, 3, data, [Inches(3.5), Inches(2.5), Inches(6.3)], top=Inches(1.3))
    bullets = [
        "USDC on Base, 암호화폐 보유: ~1% + $0.01 (최소)",
        "USDC -> USD 환전 -> 은행: ~2-2.5%",
        "Coinbase 유저 간 USDC on Base: ~$0.01 (무료)",
        "참고: Stripe 카드 2.9% + $0.30, PayPal 2.49% + $0.49",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.2), bullets, font_size=12)

    # SLIDE 32
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 정산 시나리오", 32)
    data = [
        ["항목", "Self-Managed (비수탁형)", "Coinbase-Managed (수탁형)"],
        ["커스터디", "자체 (12단어 시드 구문)", "Coinbase Exchange 연동"],
        ["정산 통화", "암호화폐 전용", "암호화폐 + USD"],
        ["법정화폐 인출", "불가", "가능"],
        ["키 관리", "가맹점 책임", "Coinbase 관리"],
    ]
    add_styled_table(slide, 5, 3, data, [Inches(2.5), Inches(5.0), Inches(5.0)], top=Inches(1.3))
    data2 = [
        ["조건", "속도"],
        ["Coinbase 유저 결제", "즉시"],
        ["Base 네트워크", "~200ms (실사용 ~2초)"],
        ["Protocol 캡처", "캡처 요청 시 즉시"],
        ["법정화폐 인출", "1-3 영업일"],
    ]
    add_styled_table(slide, 5, 2, data2, [Inches(4.0), Inches(8.3)], top=Inches(4.2))

    # SLIDE 33
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 환불 시나리오", 33)
    bullets = [
        "Self-Managed: 가맹점 직접 온체인 전송 (API 자동 환불 불가, 대시보드 환불 버튼 없음)",
        "Coinbase-Managed: Exchange 인터페이스에서 처리 (법정화폐 재환전 필요)",
        "Protocol: Void(캡처 전) / Refund(refundExpiry 전) / Reclaim(authExpiry 후) / 환불 불가(refundExpiry 후)",
        "환불 수수료: 별도 없음, 가스비 가맹점 부담, 기존 1% 수수료 미반환(추정)",
        "분쟁: 45 영업일 내 검토, 차지백 없음",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.5), bullets, font_size=12)
    add_shape(slide, Inches(9.0), Inches(5.5), Inches(3.8), Inches(0.6), fill_hex=COLOR_NEGATIVE)
    add_textbox(slide, Inches(9.0), Inches(5.5), Inches(3.8), Inches(0.6),
                "환불 자동화: 업계 최하위", font_size=12, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # SLIDE 34
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 비즈니스 모델 및 기술 아키텍처", 34)
    data = [
        ["구분", "설명"],
        ["핵심 수익원", "거래 수수료 1%"],
        ["2차", "USDC 준비금 이자 (Circle 50:50)"],
        ["3차", "환전 스프레드 0.5-2%"],
        ["간접", "생태계 락인 (Commerce -> Exchange/Custody)"],
        ["전략적", "Base 시퀀서 수수료"],
    ]
    add_styled_table(slide, 6, 2, data, [Inches(3.0), Inches(9.3)], top=Inches(1.3))
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.2),
                "전사: 연 매출 $7.181B, USDC 수익 ~$1.35B (19%), 시가총액 ~$46B\n"
                "Commerce = 직접 수익보다 'USDC 생태계 확장 + 가맹점 락인' 전략적 가치\n"
                "기술: AWS, EKS/EC2, Kafka, PostgreSQL/MongoDB, Go/Ruby/Python/React\n"
                "설계 원칙: 원자적 결제, 지갑 비종속, 불변 결제 의도, 신뢰 최소화",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 35
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] 사용자 인사이트 -- 강점 & 약점", 35)
    add_two_column(slide,
                   "가맹점 Pain Points (심각도 x 빈도)",
                   ["[매우높음x매우높음] 고객 지원 부재/지연", "[높음x높음] 환불 프로세스 복잡/수동", "[높음x매우높음] 플랫폼 전환 불확실성", "[높음x높음] 법정화폐 정산 제한", "[높음x중간] API/플러그인 방치"],
                   "소비자 Pain Points & 평점",
                   ["결제 옵션 축소", "수수료 불투명", "결제 보류 시 지원 부재", "G2 4.0/5, Capterra 4.4/5", "Trustpilot 1.3/5 (극단적 괴리)"],
                   left_color=COLOR_NEGATIVE, right_color=COLOR_NEUTRAL)

    # SLIDE 36
    slide = create_base_slide(prs, SECTION_PART1C, "[Coinbase Commerce] SWOT 분석", 36)
    add_swot_slide(slide,
                   ["Base L2 소유 (~$0.01, ~200ms)", "Commerce Payments Protocol (유일한 오픈소스)", "Shopify 독점 파트너십, 1억+ 사용자", "USDC 전략적 포지션, 비용 우위"],
                   ["환불 자동화 부재 (업계 최하위)", "고객 지원 부재", "법정화폐 정산 제한 (USD만, Managed만)", "플랫폼 불안정 (Commerce 종료)"],
                   ["GENIUS Act 수혜", "Shopify 미개척 (0.1% 미만)", "B2B 국경 간 결제", "Coinbase Business 통합"],
                   ["PayPal/Stripe 진입", "Circle CPN 장기 위협", "BitPay 38개국 정산 우위", "규제 리스크 (MiCA)"])

    # ==================================================================
    # PART 1-D: STRIPE CRYPTO (Slides 37-48)
    # ==================================================================

    # SLIDE 37
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 전략 전체상 및 수직 통합", 37)
    add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.35),
                '"가맹점은 크립토를 전혀 모르면서도, 크립토의 저비용/글로벌 장점을 누릴 수 있는 인프라"',
                font_size=12, color=ACCENT_STRIPE, bold=True)
    stack_layers = [
        ("#3A2F6E", "발행 (Issuance)", "Bridge Open Issuance"),
        ("#352B65", "블록체인 (Infrastructure)", "Tempo (Stripe + Paradigm)"),
        ("#30275C", "월렛 (Wallet)", "Privy (2025.06 인수)"),
        ("#2B2353", "결제 (Payment)", "Stripe 코어"),
        ("#261F4A", "컴플라이언스", "Bridge OCC 면허"),
        ("#211B41", "자산 관리 (Reserve)", "BlackRock, Fidelity"),
    ]
    for i, (bg, role, base) in enumerate(stack_layers):
        y = 1.65 + i * 0.75
        shape = add_shape(slide, Inches(0.5), Inches(y), Inches(9.0), Inches(0.65), fill_hex=bg)
        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, f"  {role}  |  {base}", font_size=11, color=TEXT_PRIMARY)

    add_rounded_rect(slide, Inches(10.0), Inches(2.0), Inches(2.8), Inches(2.0),
                      fill_hex=BG_CARD, border_hex=ACCENT_STRIPE, border_width=Pt(2))
    add_textbox(slide, Inches(10.1), Inches(2.1), Inches(2.6), Inches(1.8),
                "구조적 해자\n\n아무도 이 수준의 수직 통합 미보유\n\nBridge 인수 $1.1B",
                font_size=10, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

    # SLIDE 38
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 제품 포트폴리오 및 타임라인", 38)
    data = [
        ["제품", "론칭", "핵심 기능"],
        ["Stablecoin Payments", "2025.12", "스테이블코인 결제, USD 자동 정산"],
        ["Financial Accounts", "2025.05", "101개국, 잔고 보유/송수신"],
        ["Open Issuance (Bridge)", "2025.09", "자체 스테이블코인 발행"],
        ["Crypto Onramp", "2022.12~", "법정화폐 -> 크립토 전환"],
        ["Pay with Crypto", "2026.01", "크립토 잔고로 직접 결제"],
        ["x402 Protocol", "2025~", "HTTP 402 머신 결제"],
        ["MPP", "2026.03", "Tempo 기반 에이전트 결제"],
        ["Tempo 블록체인", "2026.03", "결제 전용 L1, 100K+ TPS"],
    ]
    add_styled_table(slide, 9, 3, data, [Inches(3.5), Inches(2.0), Inches(6.8)], top=Inches(1.3), font_size=9)
    # Timeline at bottom
    tl_items = ["2024.10\nBridge 인수", "2025.05\nFin.Accounts", "2025.09\nIssuance", "2025.12\nPayments", "2026.01\nPay w/Crypto", "2026.03\nTempo+MPP"]
    add_shape(slide, Inches(1.0), Inches(6.0), Inches(11.0), Inches(0.02), fill_hex=TEXT_TERTIARY)
    for i, ms in enumerate(tl_items):
        x = 1.0 + i * 2.2
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(5.9), Inches(0.18), Inches(0.18))
        oval.fill.solid()
        oval.fill.fore_color.rgb = hex_to_rgb(ACCENT_STRIPE)
        oval.line.fill.background()
        add_textbox(slide, Inches(x - 0.5), Inches(6.15), Inches(1.5), Inches(0.6),
                    ms, font_size=7, color=TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

    # SLIDE 39
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] Stablecoin Payments", 39)
    steps = ["POST /v1/\npayment_intents", "confirmPayment()\n리디렉션", "월렛 연결\n토큰/체인 선택", "온체인 제출\n블록 확인", "succeeded\n웹훅"]
    add_flow_steps(slide, steps, 1.4, ACCENT_STRIPE, box_w=2.2, box_h=0.7)
    data = [
        ["항목", "상세"],
        ["지원 스테이블코인", "USDC(ETH/SOL/Polygon/Base), USDP, USDG"],
        ["지원 월렛", "400개 이상"],
        ["수수료", "1.5% (가스비 포함)"],
        ["가맹점 요건", "미국 비즈니스만"],
        ["구독 결제", "스마트 컨트랙트 기반 -- 업계 최초"],
    ]
    add_styled_table(slide, 6, 2, data, [Inches(3.0), Inches(9.3)], top=Inches(3.5))

    # SLIDE 40
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] Pay with Crypto & Crypto Onramp", 40)
    add_card(slide, 0.5, 1.3, 6.0, 5.3, ACCENT_STRIPE, "Pay with Crypto (Crypto.com 연동)",
             ["Crypto.com 앱 내 크립토 잔고로 Stripe 가맹점 결제",
              "법정화폐 자동 전환, 사전 환전 불필요",
              "차지백 없음 (가맹점 유리)",
              "미국 우선, 글로벌 확대 예정"])
    add_card(slide, 6.8, 1.3, 6.0, 5.3, SECTION_PART2, "Crypto Onramp",
             ["법정화폐 -> 크립토 전환 (Onramp만)",
              "임베디드 위젯 또는 Stripe-hosted",
              "USDC(5개 체인), ETH, MATIC, AVAX, XLM",
              "수수료: ~5%(카드), ~1.5%(ACH)",
              "내장 KYC + Stripe Radar",
              "미국(하와이 제외) 및 EU"])

    # SLIDE 41
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] x402 & MPP -- AI 에이전트 결제", 41)
    add_card(slide, 0.5, 1.3, 6.0, 2.5, ACCENT_BASEPAY, "x402 (Coinbase 개발)",
             ["HTTP 402 기반, USDC on Base, 건별 즉시 정산",
              "최근 30일: 7,500만+ 트랜잭션, 2,400만+ 달러"])
    add_card(slide, 6.8, 1.3, 6.0, 2.5, ACCENT_STRIPE, "MPP (Stripe/Tempo, 2026.03)",
             ["Tempo 기반, 멀티 결제수단",
              "SPTs -- 에이전트 간 결제 권한 위임",
              "배치 정산, 서브-100ms"])
    data = [
        ["구분", "x402", "MPP"],
        ["인프라", "Base L2 (Coinbase)", "Tempo L1 (Stripe/Paradigm)"],
        ["결제 수단", "USDC만", "멀티 (스테이블코인+카드+BTC)"],
        ["정산", "건별 즉시", "배치 정산"],
        ["지연시간", "~2초", "서브-100ms"],
        ["유스케이스", "단건 API 호출", "대량 마이크로페이먼트"],
    ]
    add_styled_table(slide, 6, 3, data, [Inches(2.5), Inches(5.0), Inches(5.0)], top=Inches(4.1))

    # SLIDE 42
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 정산 시나리오", 42)
    steps = ["온체인 수취\n(0.4~12초)", "Bridge\n오케스트레이션", "Stripe 잔고 반영\n(카드와 통합)", "은행 출금\n(T+2 영업일)"]
    add_flow_steps(slide, steps, 1.8, ACCENT_STRIPE, box_w=2.8, box_h=0.8)
    bullets = [
        "USDC -> USD 자동 정산 (full shielding): 가맹점은 크립토를 전혀 보유/관리하지 않음",
        "USDC 직접 수취: Stablecoin Financial Accounts (101개국, 다중 통화, 가상/실물 카드)",
        "정산 수수료 1.5%에 포함",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(2.0), bullets, font_size=12)
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(0.04), Inches(0.5), fill_hex=COLOR_NEUTRAL)
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.5),
                "사용자 반응: '온체인 즉시 확정인데 정산 T+2, 스테이블코인 장점 상쇄'",
                font_size=11, color=COLOR_NEUTRAL)

    # SLIDE 43
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 수수료 구조 종합", 43)
    data = [
        ["제품", "수수료", "비고"],
        ["Stablecoin Payments", "1.5%", "USD 정산 포함, 가스비 포함"],
        ["Crypto Onramp", "~5%(카드), ~1.5%(ACH)", "결제 수단별 변동"],
        ["Pay with Crypto", "미공개", "표준 Stripe 수수료 추정"],
        ["x402", "미공개 (프리뷰)", "프로토콜 제로, Stripe 정산 별도"],
        ["MPP", "거의 제로", "배치 정산"],
    ]
    add_styled_table(slide, 6, 3, data, [Inches(3.5), Inches(3.5), Inches(5.3)], top=Inches(1.3))
    # Cost comparison bars
    cost_data = [
        ("Stripe Stablecoin", "$150", 150, ACCENT_STRIPE),
        ("Stripe 카드", "$320", 320, SECTION_PART0),
        ("Coinbase", "~$100+", 100, ACCENT_COINBASE),
        ("BitPay", "~$125", 125, SECTION_PART0),
        ("Binance Pay", "~$180+", 180, ACCENT_BINANCE),
    ]
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(3.0), Inches(0.3),
                "$10,000 월간 거래량 총비용 비교:", font_size=11, color=TEXT_PRIMARY, bold=True)
    for i, (name, label, val, color) in enumerate(cost_data):
        y = 4.9 + i * 0.38
        bar_w = val / 320 * 8.0  # Max width proportional
        add_shape(slide, Inches(3.5), Inches(y), Inches(bar_w), Inches(0.3), fill_hex=color)
        add_textbox(slide, Inches(0.5), Inches(y), Inches(2.8), Inches(0.3),
                    name, font_size=9, color=TEXT_SECONDARY)
        add_textbox(slide, Inches(3.5 + bar_w + 0.1), Inches(y), Inches(1.5), Inches(0.3),
                    label, font_size=9, color=TEXT_PRIMARY, bold=True)

    # SLIDE 44
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 환불 시나리오", 44)
    steps = ["가맹점 API\nPOST /v1/refunds", "Stripe 가맹점\n잔고(USD) 차감", "USD -> 원래\n스테이블코인 전환", "고객 원래 월렛\n온체인 전송"]
    add_flow_steps(slide, steps, 1.4, ACCENT_STRIPE, box_w=2.8, box_h=0.7)
    data = [
        ["항목", "상세"],
        ["분쟁/차지백", "미지원 (Stripe Disputes 크립토 미적용)"],
        ["잘못된 주소 송금", "복구 불가"],
        ["업계 전체 한계", "모든 크립토 결제 서비스가 전통적 차지백 미지원"],
    ]
    add_styled_table(slide, 4, 2, data, [Inches(3.0), Inches(9.3)], top=Inches(3.5))
    add_rounded_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.7),
                      fill_hex=BG_CARD, border_hex=COLOR_OPPORTUNITY, border_width=Pt(2))
    add_textbox(slide, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.6),
                "경쟁 공백: 스마트 컨트랙트 에스크로 + 분쟁 중재를 최초 구현하는 플레이어가 표준 선점",
                font_size=11, color=COLOR_OPPORTUNITY, bold=True)

    # SLIDE 45
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] Tempo 블록체인 핵심 사양", 45)
    data = [
        ["항목", "상세"],
        ["유형", "결제 전용 L1, EVM 호환"],
        ["실행 엔진", "Reth (Paradigm, Rust 기반)"],
        ["합의", "Simplex BFT (Commonware)"],
        ["TPS", "100,000+ (목표 1M)"],
        ["최종 확정성", "~0.6초, 결정론적 (리오그 없음)"],
        ["가스비", "스테이블코인으로 납부"],
        ["네이티브 토큰", "없음 (규제 리스크 회피)"],
        ["검증자", "Stripe, Visa, Zodia (허가형)"],
        ["프라이버시", "옵트인 (트랜잭션 은닉 가능)"],
    ]
    add_styled_table(slide, 10, 2, data, [Inches(3.0), Inches(9.3)], top=Inches(1.3), font_size=10,
                      row_height=Inches(0.34))
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.7),
                "파트너: Visa(검증자+카드 100개국+), Zodia(기관 커스터디), Paradigm(공동 인큐베이터)\n"
                "커뮤니티: '아무도 새 체인 불필요' 중앙화 우려 vs '실제 결제 트래픽 확보'",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 46
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 비즈니스 모델 및 파트너십", 46)
    data = [
        ["유형", "해당 제품", "신뢰도"],
        ["트랜잭션 1.5%", "Stablecoin Payments, x402, Pay w/Crypto", "확인"],
        ["전환 ~5%", "Crypto Onramp, Financial Accounts", "확인"],
        ["인프라 이용료", "Open Issuance, Bridge API", "추정"],
        ["준비금 이자", "USDB (BlackRock 등)", "추정"],
        ["네트워크 수수료", "Tempo", "추정"],
        ["생태계 교차 판매", "Radar, Billing, Connect", "확인"],
        ["FX 스프레드", "스테이블코인-법정화폐 전환", "추정"],
    ]
    add_styled_table(slide, 8, 3, data, [Inches(2.5), Inches(5.5), Inches(4.3)], top=Inches(1.3), font_size=10,
                      row_height=Inches(0.34))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.0),
                "핵심 파트너십: Visa, Zodia, Crypto.com, BlackRock/Fidelity, Circle, Coinbase, Cloudflare/Google, Paradigm, Payoneer",
                font_size=12, color=TEXT_SECONDARY)

    # SLIDE 47
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] 사용자 인사이트 -- 강점 & 약점", 47)
    add_two_column(slide,
                   "긍정 (가맹점 55%, 소비자 40%)",
                   ['"crypto 파라미터 추가만으로 수취"', 'Shadeform: 국제 카드 4.5%->1.5% (66% 절감)', '10만+ 가맹점 온보딩', '카드 대비 12% 높은 전환율', '400+ 월렛, QR 간편'],
                   "부정 (가맹점 20%, 소비자 25%)",
                   ['미국 전용 (가장 빈번)', '1.5% 수수료 논쟁', 'T+2 정산', 'USD 단일 정산', 'Tempo 반발 강함 (크립토 커뮤니티)'])
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "개선 요청 Top 5: 1)글로벌 확대 2)수수료 인하 3)코인 확대 4)다통화 정산 5)T+0 정산",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 48
    slide = create_base_slide(prs, SECTION_PART1D, "[Stripe Crypto] SWOT 분석", 48)
    add_swot_slide(slide,
                   ["기존 가맹점 네트워크 락인", "법정화폐 완전 차폐 (zero-crypto)", "수직 통합 (Bridge+Tempo+Privy)", "에이전트 결제 3중 전략", "구독 결제 스마트 컨트랙트 (업계 최초)"],
                   ["미국 전용 수취", "1.5% 프리미엄", "T+2 정산, 분쟁 미지원", "스테이블코인 위주, USD 단일 정산"],
                   ["글로벌 확장 (101개국 기반)", "AI 에이전트 결제 선점", "Tempo T+0 정산", "B2B 크로스보더 ($4,000억)"],
                   ["수수료 경쟁 (Coinbase 1%, Solana ~0%)", "Tempo 중앙화 반발", "x402 채택 부진", "PayPal 4.3억 계정"])

    # ==================================================================
    # PART 2: CROSS-COMPARISON (Slides 49-56)
    # ==================================================================

    # SLIDE 49
    slide = create_base_slide(prs, SECTION_PART2, "수수료 비교 종합", 49)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["결제 수수료", "1%", "1%", "1% (외부), 0% (CB유저)", "1.5%"],
        ["네트워크 가스비", "~$0.01 (구매자)", "0% (오프체인)", "~$0.01 (Base)", "포함"],
        ["DEX 스왑", "변동 (구매자)", "FX 0.1~0.5%", "~0.3%", "포함"],
        ["Payout 수수료", "N/A", "0.80%", "N/A", "N/A"],
        ["법정화폐 전환", "+1-5 영업일", "수동(0.1%+출금비)", "~1-1.5%", "포함 (0%)"],
        ["$10K 법정화폐 총비용", "~$100+ (변동)", "~$140~150", "~$100+ (변동)", "$150"],
        ["실효 수수료율", "~1%/~2%+", "~1.4~1.5%", "~1%/~2-2.5%", "1.5% (all-in)"],
    ]
    add_comparison_table(slide, headers, data)
    add_source(slide, "추가: BitPay 1%+$0.25, PayPal 0.99%/1.5%, NOWPayments 0.5%, Solana Pay ~0%, BTCPay 0%")

    # SLIDE 50
    slide = create_base_slide(prs, SECTION_PART2, "정산 비교 종합", 50)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["정산 속도", "~2초", "즉시 (10ms)", "~200ms (Base)", "0.4~12초"],
        ["법정화폐 정산", "간접 (CB 경유)", "불가 (수동)", "간접 (Managed, USD)", "직접 (자동 USD)"],
        ["법정화폐 출금", "+1-5 영업일", "수동 처리", "+1-3 영업일", "T+2 영업일"],
        ["정산 통화", "USDC, ETH, 기타", "USDT 기본 (50종+)", "USDC, ETH, USD", "USD"],
        ["에스크로", "온체인 지원", "미지원", "온체인 지원", "오프체인 지원"],
        ["부분 캡처", "지원", "미지원", "지원 (Protocol)", "지원"],
        ["가맹점 크립토 노출", "있음", "있음", "있음", "없음 (완전 차폐)"],
    ]
    add_comparison_table(slide, headers, data)

    # SLIDE 51
    slide = create_base_slide(prs, SECTION_PART2, "환불 비교 종합", 51)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["자동 환불", "지원 (2세대)", "API + 대시보드", "제한적", "API + 대시보드"],
        ["부분 환불", "지원", "지원", "지원 (Protocol)", "지원"],
        ["Void", "지원", "미지원", "지원", "N/A"],
        ["Reclaim", "지원", "미지원", "지원", "미지원"],
        ["환불 통화", "원래 결제 토큰", "정산 통화(USDT)", "USDC", "원래 스테이블코인"],
        ["차지백", "없음", "없음", "없음", "없음"],
        ["구매자 보호", "없음", "없음", "에스크로 부분", "없음"],
        ["환불 UX 수준", "중 (개선 중)", "중", "하 (업계 최하위)", "상"],
    ]
    add_comparison_table(slide, headers, data)

    # SLIDE 52
    slide = create_base_slide(prs, SECTION_PART2, "기능 매트릭스 종합", 52)
    headers = ["기능", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["지원 코인", "Uniswap V3 전체", "50종+/300종+(P2P)", "100+종", "스테이블코인 3종"],
        ["P2P 송금", "별도", "무료, 즉시", "별도", "미지원"],
        ["구독 결제", "SpendPermission", "미지원", "지원 (온체인)", "스마트 컨트랙트"],
        ["오프라인 결제", "Flexa 통합", "QR (오프라인)", "제한적", "미지원"],
        ["오픈소스", "프로토콜", "아니오", "프로토콜", "아니오"],
        ["Shopify 통합", "네이티브", "제한적", "네이티브", "기본 내장"],
        ["AI 에이전트 결제", "x402", "미지원", "미지원", "x402+MPP+ACP"],
        ["커스터디", "비수탁+수탁", "수탁", "비수탁+수탁", "수탁"],
        ["통합 난이도", "중간", "낮음~높음", "중간", "매우 낮음"],
        ["개발자 SDK", "JS 중심", "iOS/Android만", "오픈소스", "Stripe SDK 전체"],
    ]
    cw = [Inches(2.0), Inches(2.6), Inches(2.6), Inches(2.6), Inches(2.5)]
    add_comparison_table(slide, headers, data, col_widths=cw, font_size=9, header_font_size=9)

    # SLIDE 53
    slide = create_base_slide(prs, SECTION_PART2, "경쟁력 스코어카드 (10점 만점)", 53)
    headers = ["평가 항목", "Base Pay", "Binance Pay", "Coinbase", "Stripe", "PayPal"]
    data = [
        ["수수료 경쟁력", "7", "7", "7", "6", "7"],
        ["정산 속도", "9", "10", "10", "7", "8"],
        ["법정화폐 정산", "4", "2", "5", "10", "10"],
        ["환불/분쟁", "7", "6", "5", "7", "10"],
        ["지원 암호화폐", "8", "10", "8", "2", "8"],
        ["통합 용이성", "6", "6", "7", "10", "10"],
        ["사용자 기반", "7", "10", "7", "8", "10"],
        ["규제 안정성", "8", "5", "8", "10", "10"],
        ["종합", "7.0", "7.0", "7.3", "7.5", "9.1"],
    ]
    cw6 = [Inches(2.0)] * 6
    hdr_colors = [BG_TABLE_HEADER, COMP_HDR_BASEPAY, COMP_HDR_BINANCE, COMP_HDR_COINBASE, COMP_HDR_STRIPE, "#1A3050"]
    add_comparison_table(slide, headers, data, col_widths=cw6, service_header_colors=hdr_colors, font_size=10)

    # SLIDE 54
    slide = create_base_slide(prs, SECTION_PART2, "비즈니스 모델 비교", 54)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["핵심 수익원", "시퀀서비+USDC+1%", "MDR 1% ~$110M", "거래 수수료 1%+USDC", "트랜잭션 1.5%+전환 ~5%"],
        ["전체 내 비중", "생태계 인프라", "0.6% (전략 도구)", "소규모 (생태계 락인)", "크립토는 확장 중"],
        ["전략적 역할", "Base 트래픽+USDC 유통", "사용자 활성화+자금 체류", "USDC 생태계+가맹점 락인", "기존 가맹점 전환"],
        ["오픈소스 전략", "Apache 2.0", "폐쇄적", "프로토콜 오픈소스", "폐쇄적"],
    ]
    add_comparison_table(slide, headers, data)

    # SLIDE 55
    slide = create_base_slide(prs, SECTION_PART2, "기술 아키텍처 비교", 55)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["결제 레이어", "온체인 (스마트 컨트랙트)", "오프체인 (내부 원장)", "온체인 (스마트 컨트랙트)", "온+오프체인 하이브리드"],
        ["블록체인", "Base L2 (OP Stack)", "N/A", "Base L2", "ETH/SOL/Polygon+Tempo"],
        ["성능", "~2초 정산", "10,000+ TPS, 10ms", "~200ms (Base)", "Tempo: 100K+ TPS"],
        ["DEX 통합", "Uniswap V3", "내부 FX Engine", "Uniswap V3", "Bridge 오케스트레이션"],
        ["보안", "Spearbit 감사", "HMAC-SHA512, KYC 3단계", "Spearbit 감사", "Stripe Radar, OCC"],
        ["월렛 지원", "모든 EVM 지갑", "Binance 앱만", "모든 EVM 지갑", "400+ 월렛 (Privy)"],
    ]
    add_comparison_table(slide, headers, data)

    # SLIDE 56
    slide = create_base_slide(prs, SECTION_PART2, "사용자 평점 및 핵심 Pain Points 비교", 56)
    headers = ["항목", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["Capterra/G2", "4.4/5", "4.5/5", "4.0~4.4/5", "4.0/5"],
        ["Trustpilot", "N/A (BBB F)", "1.4/5", "1.3/5", "2.7/5"],
        ["최대 불만 1", "고객 지원 부재", "고객 지원 품질", "고객 지원 부재", "미국 전용 제한"],
        ["최대 불만 2", "Commerce 종료", "계정 동결", "환불 수동", "1.5% 수수료"],
        ["최대 불만 3", "법정화폐 정산", "Binance 계정 필수", "플랫폼 전환 불확실", "T+2 정산"],
    ]
    add_comparison_table(slide, headers, data)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8),
                "공통 Pain Point: 고객 지원 품질 (4개 중 3개), 법정화폐 정산 제한 (3개)\n"
                "차별적: Stripe만 고객 지원 불만이 상위가 아님 (대신 지역 제한)",
                font_size=12, color=TEXT_PRIMARY, bold=True)

    # ==================================================================
    # PART 3: STRATEGIC IMPLICATIONS (Slides 57-64)
    # ==================================================================

    # SLIDE 57
    slide = create_base_slide(prs, SECTION_PART3, 'AI 에이전트 결제 경쟁 -- "에이전트 결제 전쟁"', 57)
    headers = ["프로토콜", "주도", "결제 모델", "4개 서비스 관계"]
    data = [
        ["x402", "Coinbase", "요청당 즉시 (HTTP 402)", "Base Pay(주도), Stripe(멤버)"],
        ["MPP", "Stripe/Tempo", "세션 배치 정산", "Stripe(주도)"],
        ["ACP", "OpenAI/Stripe", "에이전트-가맹점 체크아웃", "Stripe(인프라)"],
        ["AP2", "Google", "지출 권한 관리", "간접"],
        ["TAP", "Visa", "카드 네트워크 위 에이전트", "Stripe(Tempo 검증자)"],
        ["Agent Ready", "PayPal", "PayPal 생태계", "경쟁"],
    ]
    cw = [Inches(2.0), Inches(2.5), Inches(3.5), Inches(4.3)]
    add_styled_table(slide, 7, 4, [headers] + data, cw, top=Inches(1.3), font_size=10)
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0),
                "Stripe 이중 전략: x402+MPP+ACP 3개 동시 참여 -- 어떤 표준이든 Stripe 존재\n"
                "Base Pay x402: 1.19억 건, 연간화 $600M (다만 실질 일일 ~$28K)\n"
                "Binance Pay / Coinbase Commerce: AI 에이전트 결제 직접 참여 없음",
                font_size=11, color=TEXT_SECONDARY)

    # SLIDE 58
    slide = create_base_slide(prs, SECTION_PART3, "시장 트렌드와 4개 서비스의 대응 전략", 58)
    headers = ["트렌드", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["스테이블코인 결제", "USDC 기본, x402 98.7%", "B2C 98% 스테이블코인", "USDC 자동 환전", "USDC/USDP/USDG"],
        ["L2 저비용 결제", "Base L2 소유 (~$0.01)", "오프체인 (0원)", "Base L2 활용", "+Tempo 자체 L1"],
        ["AI 에이전트 결제", "x402 주도", "미참여", "미참여", "x402+MPP+ACP 3중"],
        ["전통 결제사 진출", "오픈소스로 대응", "독자 생태계", "Shopify 파트너십", "자체가 전통 결제사"],
        ["규제 명확화", "미국 상장사, MSB", "70개국+ 리스크", "GENIUS Act 수혜", "OCC 면허, 가장 유리"],
    ]
    add_comparison_table(slide, headers, data)

    # SLIDE 59
    slide = create_base_slide(prs, SECTION_PART3, "각 서비스 권고사항 종합", 59)
    headers = ["우선순위", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["P0", "24/7 고객 지원 구축", "고객 지원 전면 개편", "실시간 지원 채널 구축", "가맹점 수취 글로벌 확대"],
        ["P0", "법정화폐 직접 정산", "법정화폐 직접 정산", "환불 자동화", "볼륨 할인 티어 (1%이하)"],
        ["P1", "온체인 구매자 보호", "구독/반복 결제", "총비용 투명화", "Tempo T+0 정산"],
        ["P1", "자동 환불 워크플로우", "이커머스 플러그인 확대", "다국통화 법정화폐", "온체인 분쟁 해결 파일럿"],
        ["P2", "다중 언어 SDK", "서버사이드 SDK+문서", "Protocol 멀티 플랫폼", "다통화 정산 (EUR, GBP)"],
        ["P2", "x402 스펙-구현 일관성", "게스트 체크아웃", "API 문서+플러그인 유지보수", "BTC/ETH 직접 수취"],
    ]
    cw = [Inches(1.5), Inches(2.7), Inches(2.7), Inches(2.7), Inches(2.7)]
    hdr_colors = [BG_TABLE_HEADER, COMP_HDR_BASEPAY, COMP_HDR_BINANCE, COMP_HDR_COINBASE, COMP_HDR_STRIPE]
    add_comparison_table(slide, headers, data, col_widths=cw, service_header_colors=hdr_colors, font_size=9)

    # SLIDE 60
    slide = create_base_slide(prs, SECTION_PART3, "전략적 시사점 -- 5대 핵심 발견", 60)
    insights = [
        ("법정화폐 정산이 메인스트림 채택의 최대 결정 요인",
         "4개 서비스 중 자동 법정화폐 정산은 Stripe만 제공. 전통 사업자 온보딩의 핵심 차별화 요소"),
        ('"독립 수익 사업"은 없다 -- 모두 전략적 인프라',
         "Base Pay(시퀀서비+USDC), Binance Pay(0.6%), Coinbase(락인), Stripe(전환). 결제 자체 수익보다 생태계 가치가 핵심"),
        ("환불/소비자 보호가 업계 전체의 미해결 과제",
         "4개 서비스 모두 차지백/구매자 보호 미제공. Base Pay 에스크로가 가장 진전된 시도"),
        ("Stripe의 수직 통합이 구조적 위협",
         "Bridge+Tempo+Privy+Stripe 풀스택은 복제 불가. 미국 전용+1.5%가 단기 기회 창"),
        ("AI 에이전트 결제가 차세대 성장 동인",
         "Stripe 3개 동시 참여로 가장 유리. Binance Pay/Coinbase Commerce는 부재"),
    ]
    add_numbered_blocks(slide, insights, SECTION_PART3, y_start=1.3, block_h=1.05, gap=0.1)

    # SLIDE 61
    slide = create_base_slide(prs, SECTION_PART3, "서비스별 최적 사용 시나리오", 61)
    data = [
        ["시나리오", "최적 서비스", "근거"],
        ["실물 이커머스 (배송 후 결제)", "Base Pay / Coinbase", "에스크로 (Auth-Capture) 지원"],
        ["디지털 상품 즉시 결제", "Base Pay, Coinbase", "즉시 정산, 저비용"],
        ["법정화폐 정산 필수", "Stripe Crypto", "자동 USD 정산, 기존 Stripe 통합"],
        ["Binance 사용자 기반", "Binance Pay", "P2P 무료, 3억 사용자"],
        ["글로벌 법정화폐 정산", "BitPay (4개 외)", "38개국 직접 정산"],
        ["AI 에이전트 API 과금", "Base Pay / Stripe", "건당 마이크로페이먼트"],
        ["SaaS 구독 결제", "Stripe Crypto", "스마트 컨트랙트 반복 결제"],
        ["신흥국/금융 포용", "Binance Pay", "사용자 기반, 오프체인 무가스비"],
        ["크립토 네이티브 / 비용 최소화", "Base Pay / Coinbase", "1% + $0.01, 오픈소스"],
        ["고가 물품 (차지백 회피)", "모든 크립토 결제", "비가역적 결제로 사기 분쟁 방지"],
    ]
    add_styled_table(slide, 11, 3, data, [Inches(3.5), Inches(3.5), Inches(5.3)], top=Inches(1.3),
                      font_size=10, row_height=Inches(0.34))

    # SLIDE 62
    slide = create_base_slide(prs, SECTION_PART3, "데이터 신뢰도 평가 종합", 62)
    headers = ["분석 영역", "Base Pay", "Binance Pay", "Coinbase Commerce", "Stripe Crypto"]
    data = [
        ["결제 프로세스", "높음 (오픈소스)", "높음 (공식 문서)", "높음 (오픈소스)", "높음 (공식 문서)"],
        ["정산 프로세스", "높음", "높음", "높음", "높음"],
        ["환불/분쟁", "중-높", "높음", "중-높", "높음"],
        ["수수료", "높음", "높음 (FX 비공개)", "중-높", "높음 (일부 미공개)"],
        ["시장 데이터", "중", "중-높", "중-상", "중"],
        ["사용자 인사이트", "중", "중", "중", "중"],
        ["기술 아키텍처", "높음", "중-높", "높음", "높음/중"],
    ]
    add_comparison_table(slide, headers, data)
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.7),
                "교차 검증: 4개 보고서 간 주요 수치 일관성 확인 완료\n"
                "주의: 시장 규모는 기관별 정의 차이로 2배+ 편차, x402 채택 수치는 테스트 거래 포함 의심",
                font_size=10, color=TEXT_TERTIARY)

    # SLIDE 63
    slide = create_base_slide(prs, SECTION_PART3, "결론 및 핵심 메시지", 63)
    # Key takeaway box
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.9), fill_hex=BG_CARD)
    add_shape(slide, Inches(0.5), Inches(1.3), Inches(0.05), Inches(0.9), fill_hex=SECTION_PART3)
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.8),
                "크립토 결제는 \"별도 카테고리\"에서 \"기존 결제 인프라의 일부\"로 전환 중",
                font_size=16, color=TEXT_PRIMARY, bold=True)
    conclusions = [
        ("스테이블코인이 결제의 표준으로 부상",
         "Binance Pay B2C 98%, x402 USDC 98.7%, 모든 서비스가 USDC/USDT 중심"),
        ("법정화폐 자동 정산이 채택의 분수령",
         "현재 Stripe만 자동 제공. Base Pay/Binance Pay/Coinbase Commerce 모두 간접/미지원"),
        ("AI 에이전트 결제가 차세대 전장",
         "Stripe(3중 전략)과 Base Pay(x402)가 선점 경쟁 중. 90일 내 6개 프로토콜 론칭"),
        ("운영 품질이 기술 혁신만큼 중요",
         "4개 서비스 중 3개에서 고객 지원이 최대 불만. 기술 우위가 운영 품질을 따라가지 못함"),
    ]
    for i, (title, desc) in enumerate(conclusions):
        y = 2.5 + i * 1.1
        add_textbox(slide, Inches(0.7), Inches(y), Inches(12.0), Inches(0.4),
                    f"{i+1}. {title}", font_size=13, color=SECTION_PART3, bold=True)
        add_textbox(slide, Inches(0.9), Inches(y + 0.4), Inches(11.8), Inches(0.5),
                    desc, font_size=12, color=TEXT_SECONDARY)

    # SLIDE 64
    slide = create_base_slide(prs, SECTION_PART3, "부록 -- 주요 출처 및 참고 자료", 64)
    sources = [
        ("공식 소스", "Stripe Docs, Coinbase Blog/Help, Binance Merchant Docs, GitHub"),
        ("시장 데이터", "Research Nester, GII Research, CoinLaw, McKinsey, PYMNTS, eMarketer"),
        ("경쟁사", "BitPay, CoinGate, NOWPayments, PayPal, Circle CPN, Crypto.com Pay"),
        ("규제", "GENIUS Act (US), MiCA (EU)"),
        ("사용자 리뷰", "Trustpilot, G2, Capterra, Reddit, Hacker News, GitHub Issues"),
    ]
    y = 1.5
    for cat, items in sources:
        add_textbox(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(0.35),
                    cat, font_size=14, color=SECTION_PART3, bold=True)
        add_textbox(slide, Inches(0.8), Inches(y + 0.35), Inches(11.9), Inches(0.35),
                    items, font_size=11, color=TEXT_TERTIARY)
        y += 0.85
    add_textbox(slide, Inches(0.6), Inches(y + 0.2), Inches(12.1), Inches(0.4),
                "분석 기준일: 2026-04-15", font_size=14, color=TEXT_PRIMARY, bold=True)

    return prs


# ============================================================
# 3. MAIN
# ============================================================

if __name__ == "__main__":
    output_path = "/Users/okyokwon/Desktop/Private_Projects/research_p/_workspace/crypto_payment_services_research.pptx"
    print("Building presentation...")
    prs = build_presentation()
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
